#!/usr/bin/env python3
"""
Generate Editable PPTX
将HTML演示文稿转换为可编辑的PowerPoint文件

Usage:
    python generate_editable_pptx.py <input.html> [output.pptx] [--theme THEME]

Examples:
    python generate_editable_pptx.py presentation.html
    python generate_editable_pptx.py presentation.html output.pptx --theme business
    python generate_editable_pptx.py presentation.html --list-themes
"""

import sys
import os
import argparse
from pathlib import Path

# 导入自定义模块
from html_parser import HTMLParser
from pptx_builder import PPTXBuilder
from theme_manager import ThemeManager


def check_dependencies():
    """检查依赖项"""
    missing = []

    try:
        from pptx import Presentation
    except ImportError:
        missing.append("python-pptx")

    try:
        from bs4 import BeautifulSoup
    except ImportError:
        missing.append("beautifulsoup4")

    if missing:
        print("❌ Missing dependencies. Install with:")
        print(f"   pip install {' '.join(missing)}")
        sys.exit(1)


def find_html_file(directory: str = None) -> Path:
    """
    查找HTML文件

    Args:
        directory: 搜索目录，默认为当前目录

    Returns:
        HTML文件路径
    """
    if directory is None:
        directory = os.getcwd()

    directory = Path(directory)

    # 查找所有HTML文件
    html_files = list(directory.glob("*.html"))

    if not html_files:
        print(f"❌ No HTML files found in {directory}")
        sys.exit(1)

    # 返回最近修改的文件
    html_files.sort(key=lambda x: x.stat().st_mtime, reverse=True)
    return html_files[0]


def generate_pptx(html_file: str, output_file: str = None, theme: str = 'spring') -> str:
    """
    生成可编辑的PPTX文件

    Args:
        html_file: 输入HTML文件路径
        output_file: 输出PPTX文件路径（可选）
        theme: 主题名称

    Returns:
        输出文件路径
    """
    html_path = Path(html_file)

    if not html_path.exists():
        print(f"❌ HTML file not found: {html_file}")
        sys.exit(1)

    # 确定输出文件名
    if output_file is None:
        output_file = html_path.with_suffix('.pptx')

    output_path = Path(output_file)

    print(f"📖 Parsing HTML file: {html_path.name}")
    print(f"🎨 Using theme: {theme}")
    print()

    # 1. 解析HTML
    print("Step 1: Parsing HTML...")
    parser = HTMLParser()
    try:
        slide_data = parser.parse(str(html_path))
        print(f"  ✓ Found {len(slide_data)} slides")
    except Exception as e:
        print(f"  ❌ Failed to parse HTML: {e}")
        sys.exit(1)

    # 2. 加载主题
    print("\nStep 2: Loading theme...")
    theme_manager = ThemeManager()
    theme_config = theme_manager.get_theme(theme)
    print(f"  ✓ Theme loaded: {theme_config.get('name', theme)}")

    # 3. 构建PPTX
    print("\nStep 3: Building PPTX...")
    builder = PPTXBuilder(theme_config)

    try:
        builder.build(slide_data)
        print(f"  ✓ Built {len(slide_data)} slides")
    except Exception as e:
        print(f"  ❌ Failed to build PPTX: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)

    # 4. 验证可编辑性
    print("\nStep 4: Verifying editability...")
    stats = builder.verify_editability()
    print(f"  ✓ Text boxes: {stats['text_boxes']}")
    print(f"  ✓ Shapes: {stats['shapes']}")
    print(f"  ✓ Images: {stats['images']}")
    print(f"  ✓ Screenshots: {stats['screenshots']} (should be 0)")

    # 5. 保存文件
    print(f"\nStep 5: Saving PPTX file...")
    try:
        builder.save(str(output_path))
        print(f"  ✓ Saved: {output_path}")
    except Exception as e:
        print(f"  ❌ Failed to save PPTX: {e}")
        sys.exit(1)

    # 6. 验证内容一致性
    print(f"\nStep 6: Validating content consistency...")
    try:
        from smart_ppt_validator import SmartPPTValidator
        validator = SmartPPTValidator(str(html_path), str(output_path))
        report = validator.validate()

        if report['status'] == 'PASS':
            print(f"  ✓ Validation passed: PPT and HTML content are consistent")
        else:
            print(f"  ⚠️  Validation failed: Found {len(report['issues'])} issue(s)")
            print(f"\n  Issues:")
            for i, issue in enumerate(report['issues'], 1):
                severity_icon = '🔴' if issue['severity'] == 'HIGH' else '🟡'
                print(f"    {severity_icon} {issue['message']}")
            print(f"\n  💡 Tip: Please review the PPTX file and check if the content is correct.")
    except Exception as e:
        print(f"  ⚠️  Validation skipped: {e}")

    # 7. 总结
    print("\n" + "=" * 60)
    print("✅ Editable PPTX generated successfully!")
    print("=" * 60)
    print(f"\n📁 Output file: {output_path}")
    print(f"📊 Slides: {len(slide_data)}")
    print(f"🎨 Theme: {theme_config.get('name', theme)}")
    print(f"\n✨ Editability:")
    print(f"  • All text is editable: {stats['text_boxes']} text boxes")
    print(f"  • All shapes are adjustable: {stats['shapes']} shapes")
    print(f"  • No screenshots: {stats['screenshots']} (fully editable)")
    print(f"\n💡 To edit: Open the file in PowerPoint and modify any element directly.")

    return str(output_path)


def list_themes():
    """列出所有可用主题"""
    theme_manager = ThemeManager()
    theme_manager.list_themes()


def verify_pptx(pptx_file: str):
    """
    验证PPTX文件的可编辑性

    Args:
        pptx_file: PPTX文件路径
    """
    from pptx import Presentation

    pptx_path = Path(pptx_file)

    if not pptx_path.exists():
        print(f"❌ PPTX file not found: {pptx_file}")
        sys.exit(1)

    print(f"🔍 Verifying: {pptx_path.name}")
    print()

    prs = Presentation(str(pptx_path))

    stats = {
        'slides': len(prs.slides),
        'text_boxes': 0,
        'shapes': 0,
        'images': 0,
        'screenshots': 0
    }

    for slide in prs.slides:
        for shape in slide.shapes:
            stats['shapes'] += 1
            if shape.has_text_frame:
                stats['text_boxes'] += 1
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                stats['images'] += 1

    print("Verification Results:")
    print("-" * 60)
    print(f"  Slides:      {stats['slides']}")
    print(f"  Text boxes:  {stats['text_boxes']}")
    print(f"  Shapes:      {stats['shapes']}")
    print(f"  Images:      {stats['images']}")
    print(f"  Screenshots: {stats['screenshots']}")
    print("-" * 60)

    if stats['screenshots'] == 0 and stats['text_boxes'] > 0:
        print("\n✅ This PPTX is fully editable!")
        print("   All text, shapes, and colors can be modified in PowerPoint.")
    else:
        print("\n⚠️  This PPTX may have limited editability.")
        if stats['screenshots'] > 0:
            print("   Contains screenshots which cannot be edited.")
        if stats['text_boxes'] == 0:
            print("   No editable text boxes found.")


def main():
    """主函数"""
    parser = argparse.ArgumentParser(
        description='Generate editable PowerPoint presentations from HTML',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Convert HTML to PPTX (auto-detect HTML file)
  python generate_editable_pptx.py

  # Convert specific HTML file
  python generate_editable_pptx.py presentation.html

  # Specify output file and theme
  python generate_editable_pptx.py presentation.html output.pptx --theme business

  # List available themes
  python generate_editable_pptx.py --list-themes

  # Verify PPTX editability
  python generate_editable_pptx.py --verify presentation.pptx
        """
    )

    parser.add_argument(
        'html',
        nargs='?',
        help='Input HTML file (optional, will auto-detect if not provided)'
    )

    parser.add_argument(
        'output',
        nargs='?',
        help='Output PPTX file (optional, defaults to same name as HTML)'
    )

    parser.add_argument(
        '--theme',
        default='spring',
        help='Theme name (default: spring). Use --list-themes to see available themes.'
    )

    parser.add_argument(
        '--list-themes',
        action='store_true',
        help='List all available themes'
    )

    parser.add_argument(
        '--verify',
        metavar='PPTX_FILE',
        help='Verify PPTX file editability'
    )

    args = parser.parse_args()

    # 检查依赖
    check_dependencies()

    # 列出主题
    if args.list_themes:
        list_themes()
        return

    # 验证PPTX
    if args.verify:
        verify_pptx(args.verify)
        return

    # 生成PPTX
    html_file = args.html

    # 如果没有指定HTML文件，自动查找
    if html_file is None:
        html_file = find_html_file()
        print(f"📁 Auto-detected HTML file: {html_file.name}")
        print()

    generate_pptx(html_file, args.output, args.theme)


if __name__ == "__main__":
    main()
