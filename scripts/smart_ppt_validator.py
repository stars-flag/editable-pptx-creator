"""
改进的PPT验证器
更智能地对比HTML和PPT内容，减少误报
"""

import sys
import os
from pptx import Presentation
from bs4 import BeautifulSoup


class SmartPPTValidator:
    """智能PPT验证器，对比HTML和PPT内容"""

    def __init__(self, html_file: str, pptx_file: str):
        self.html_file = html_file
        self.pptx_file = pptx_file
        self.issues = []

    def validate(self) -> dict:
        """
        验证PPT与HTML的一致性

        Returns:
            验证结果字典
        """
        # 解析HTML
        html_content = self._parse_html()

        # 解析PPT
        ppt_content = self._parse_pptx()

        # 对比内容
        self._compare_content(html_content, ppt_content)

        # 生成报告
        report = {
            'html_file': self.html_file,
            'pptx_file': self.pptx_file,
            'html_slides': len(html_content),
            'ppt_slides': len(ppt_content),
            'issues': self.issues,
            'status': 'PASS' if len(self.issues) == 0 else 'FAIL'
        }

        return report

    def _parse_html(self) -> list:
        """解析HTML文件，提取幻灯片内容"""
        with open(self.html_file, 'r', encoding='utf-8') as f:
            html_content = f.read()

        soup = BeautifulSoup(html_content, 'html.parser')
        slides = soup.find_all('section', class_='slide')

        html_slides = []

        for slide in slides:
            slide_data = {
                'title': '',
                'content': [],
                'type': 'unknown',
                'content_text': ''
            }

            # 获取标题
            h2 = slide.find('h2')
            if h2:
                slide_data['title'] = h2.get_text(strip=True)

            # 检测幻灯片类型并提取内容
            if slide.find('table'):
                slide_data['type'] = 'table'
                table = slide.find('table')
                rows = table.find_all('tr')
                for row in rows:
                    cells = [cell.get_text(strip=True) for cell in row.find_all(['th', 'td'])]
                    slide_data['content'].append(cells)
                    slide_data['content_text'] += ' '.join(cells) + ' '

            elif slide.find('div', class_='two-column'):
                slide_data['type'] = 'two_column'
                two_column = slide.find('div', class_='two-column')
                columns = two_column.find_all('div', recursive=False)
                for column in columns:
                    items = column.find_all('li')
                    column_items = [item.get_text(strip=True) for item in items]
                    slide_data['content'].append(column_items)
                    slide_data['content_text'] += ' '.join(column_items) + ' '

            elif slide.find('ul'):
                slide_data['type'] = 'list'
                items = slide.find_all('li')
                slide_data['content'] = [item.get_text(strip=True) for item in items]
                slide_data['content_text'] = ' '.join(slide_data['content'])

            elif slide.find('div', class_='card'):
                slide_data['type'] = 'cards'
                cards = slide.find_all('div', class_='card')
                for card in cards:
                    card_data = {}
                    h3 = card.find('h3')
                    if h3:
                        card_data['title'] = h3.get_text(strip=True)
                        slide_data['content_text'] += h3.get_text(strip=True) + ' '
                    paragraphs = card.find_all('p')
                    content = '\n'.join([p.get_text(strip=True) for p in paragraphs if p.get_text(strip=True)])
                    card_data['content'] = content
                    slide_data['content'].append(card_data)
                    slide_data['content_text'] += content + ' '

            html_slides.append(slide_data)

        return html_slides

    def _parse_pptx(self) -> list:
        """解析PPTX文件，提取幻灯片内容"""
        prs = Presentation(self.pptx_file)

        ppt_slides = []

        for slide in prs.slides:
            slide_data = {
                'title': '',
                'content': [],
                'type': 'unknown',
                'content_text': ''
            }

            # 提取所有文本框
            text_boxes = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        text_boxes.append({
                            'text': text,
                            'top': shape.top,
                            'height': shape.height,
                            'width': shape.width
                        })

            # 按位置排序（从上到下）
            text_boxes.sort(key=lambda x: x['top'])

            # 提取标题（通常是第一个或最大的文本框）
            if text_boxes:
                # 找到最大的文本框作为标题
                title_box = max(text_boxes, key=lambda x: x['width'] * x['height'])
                slide_data['title'] = title_box['text']
                text_boxes.remove(title_box)

            # 提取内容
            for box in text_boxes:
                text = box['text']
                slide_data['content'].append(text)
                slide_data['content_text'] += text + ' '

            # 判断幻灯片类型
            if len(slide_data['content']) > 0:
                # 检查是否是表格（包含多个短文本）
                short_texts = [t for t in slide_data['content'] if len(t) < 50]
                if len(short_texts) > 10:
                    slide_data['type'] = 'table'
                # 检查是否是列表（包含多个中等长度文本）
                elif len(slide_data['content']) > 5:
                    slide_data['type'] = 'list'
                else:
                    slide_data['type'] = 'content'

            ppt_slides.append(slide_data)

        return ppt_slides

    def _compare_content(self, html_slides: list, ppt_slides: list) -> None:
        """对比HTML和PPT内容"""
        # 检查幻灯片数量
        if len(html_slides) != len(ppt_slides):
            self.issues.append({
                'type': 'slide_count_mismatch',
                'severity': 'HIGH',
                'message': f'幻灯片数量不匹配: HTML有{len(html_slides)}页，PPT有{len(ppt_slides)}页'
            })

        # 逐页对比
        min_slides = min(len(html_slides), len(ppt_slides))
        for i in range(min_slides):
            html_slide = html_slides[i]
            ppt_slide = ppt_slides[i]

            # 检查标题
            if html_slide['title'] and ppt_slide['title']:
                if html_slide['title'] != ppt_slide['title']:
                    self.issues.append({
                        'type': 'title_mismatch',
                        'severity': 'MEDIUM',
                        'slide': i + 1,
                        'message': f'第{i+1}页标题不匹配: HTML="{html_slide["title"]}", PPT="{ppt_slide["title"]}"'
                    })

            # 检查内容文本覆盖率
            html_text = html_slide['content_text'].strip()
            ppt_text = ppt_slide['content_text'].strip()

            if len(html_text) > 0 and len(ppt_text) > 0:
                # 计算HTML内容在PPT中的覆盖率
                coverage = self._calculate_coverage(html_text, ppt_text)

                if coverage < 0.5:
                    self.issues.append({
                        'type': 'content_missing',
                        'severity': 'HIGH',
                        'slide': i + 1,
                        'message': f'第{i+1}页内容严重缺失: 覆盖率仅{coverage:.1%}，HTML长度{len(html_text)}，PPT长度{len(ppt_text)}'
                    })
                elif coverage < 0.8:
                    self.issues.append({
                        'type': 'content_partial',
                        'severity': 'MEDIUM',
                        'slide': i + 1,
                        'message': f'第{i+1}页内容部分缺失: 覆盖率{coverage:.1%}，HTML长度{len(html_text)}，PPT长度{len(ppt_text)}'
                    })

    def _calculate_coverage(self, html_text: str, ppt_text: str) -> float:
        """
        计算HTML内容在PPT中的覆盖率

        Args:
            html_text: HTML文本
            ppt_text: PPT文本

        Returns:
            覆盖率（0-1）
        """
        # 将文本分割成词
        html_words = set(html_text.split())
        ppt_words = set(ppt_text.split())

        if len(html_words) == 0:
            return 1.0

        # 计算交集
        intersection = html_words & ppt_words

        # 计算覆盖率
        coverage = len(intersection) / len(html_words)

        return coverage

    def print_report(self, report: dict) -> None:
        """打印验证报告"""
        print("\n" + "=" * 80)
        print("PPT验证报告（智能版）")
        print("=" * 80)

        print(f"\nHTML文件: {report['html_file']}")
        print(f"PPT文件: {report['pptx_file']}")
        print(f"HTML幻灯片数: {report['html_slides']}")
        print(f"PPT幻灯片数: {report['ppt_slides']}")
        print(f"验证状态: {report['status']}")

        if report['status'] == 'PASS':
            print("\n✅ 验证通过！PPT与HTML内容一致。")
        else:
            print(f"\n❌ 验证失败！发现 {len(report['issues'])} 个问题：\n")

            for i, issue in enumerate(report['issues'], 1):
                severity_icon = '🔴' if issue['severity'] == 'HIGH' else '🟡'
                print(f"{severity_icon} 问题 {i}: {issue['message']}")
                print(f"   类型: {issue['type']}")
                print(f"   严重程度: {issue['severity']}")
                if 'slide' in issue:
                    print(f"   幻灯片: 第{issue['slide']}页")
                print()

        print("=" * 80)


def main():
    """主函数"""
    if len(sys.argv) < 3:
        print("用法: python smart_ppt_validator.py <html_file> <pptx_file>")
        print("示例: python smart_ppt_validator.py 熊出没介绍.html 熊出没介绍.pptx")
        sys.exit(1)

    html_file = sys.argv[1]
    pptx_file = sys.argv[2]

    if not os.path.exists(html_file):
        print(f"错误: HTML文件不存在: {html_file}")
        sys.exit(1)

    if not os.path.exists(pptx_file):
        print(f"错误: PPTX文件不存在: {pptx_file}")
        sys.exit(1)

    # 创建验证器
    validator = SmartPPTValidator(html_file, pptx_file)

    # 执行验证
    report = validator.validate()

    # 打印报告
    validator.print_report(report)

    # 返回状态码
    sys.exit(0 if report['status'] == 'PASS' else 1)


if __name__ == '__main__':
    main()
