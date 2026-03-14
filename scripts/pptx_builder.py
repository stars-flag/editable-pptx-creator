"""
PPTX Builder Module - 修复版
使用python-pptx创建可编辑的PowerPoint演示文稿
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from typing import List, Dict, Any, Tuple
import os


class PPTXBuilder:
    """PPTX构建器，创建可编辑的PowerPoint对象"""

    def __init__(self, theme: Dict[str, Any]):
        """
        初始化PPTX构建器

        Args:
            theme: 主题配置字典
        """
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

        # 应用主题
        self.colors = self._load_theme_colors(theme.get('colors', {}))
        self.fonts = theme.get('fonts', {})
        self.layout = theme.get('layout', {})

    def _load_theme_colors(self, theme_colors: Dict[str, str]) -> Dict[str, RGBColor]:
        """加载主题颜色"""
        default_colors = {
            'primary': '#7CB342',
            'secondary': '#F48FB1',
            'accent': '#FFD54F',
            'sky': '#81D4FA',
            'text_primary': '#2E7D32',
            'text_secondary': '#5D4037',
            'text_light': '#8D6E63',
            'bg_light': '#E8F5E9',
            'bg_pink': '#FCE4EC',
        }

        colors = {}
        for key, default_hex in default_colors.items():
            hex_color = theme_colors.get(key, default_hex)
            colors[key] = self._hex_to_rgb(hex_color)

        return colors

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """将十六进制颜色转换为RGBColor对象"""
        hex_color = hex_color.lstrip('#')
        rgb = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        return RGBColor(*rgb)

    def build(self, slide_data: List[Dict[str, Any]]) -> None:
        """
        构建PPTX演示文稿

        Args:
            slide_data: 幻灯片数据列表
        """
        for data in slide_data:
            slide_type = data.get('type', 'content')
            title = data.get('title', '')

            if slide_type == 'title':
                subtitle = data.get('subtitle', '')
                self.add_title_slide(title, subtitle)

            elif slide_type == 'table':
                table = data.get('table', [])
                self.add_table_slide(title, table)

            elif slide_type == 'colors':
                colors = data.get('colors', [])
                self.add_color_cards_slide(title, colors)

            elif slide_type == 'poems':
                poems = data.get('poems', [])
                self.add_poems_slide(title, poems)

            elif slide_type == 'two_column':
                two_column = data.get('two_column', {})
                self.add_two_column_slide(title, two_column)

            elif slide_type == 'cards':
                cards = data.get('cards', [])
                self.add_cards_slide(title, cards)

            elif slide_type == 'image':
                image_data = data.get('image', {})
                self.add_image_slide(title, image_data)

            else:  # content
                items = data.get('items', [])
                self.add_content_slide(title, items)

    def add_background(self, slide) -> None:
        """添加渐变背景"""
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 135
        fill.gradient_stops[0].color.rgb = self.colors['bg_light']
        fill.gradient_stops[1].color.rgb = self.colors['bg_pink']

    def add_title_slide(self, title: str, subtitle: str = '') -> None:
        """添加标题幻灯片"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 添加背景
        self.add_background(slide)

        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(8), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(60)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['text_primary']
        title_para.font.name = self.fonts.get('heading', '微软雅黑')

        # 添加副标题
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(4), Inches(8), Inches(0.8)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.alignment = PP_ALIGN.CENTER
            subtitle_para.font.size = Pt(28)
            subtitle_para.font.color.rgb = self.colors['text_secondary']
            subtitle_para.font.name = self.fonts.get('body', '微软雅黑')

        # 添加装饰
        self.add_decoration(slide)

    def add_table_slide(self, title: str, table_data: List[List[str]]) -> None:
        """添加表格幻灯片"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 添加背景
        self.add_background(slide)

        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(8), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['text_primary']
        title_para.font.name = self.fonts.get('heading', '微软雅黑')

        if not table_data:
            return

        # 计算表格位置和大小
        num_rows = len(table_data)
        num_cols = len(table_data[0]) if table_data else 0

        # 表格起始位置
        table_x = 1.0
        table_y = 1.8
        table_width = 8.0
        row_height = 0.5
        col_width = table_width / num_cols

        # 创建表格背景
        table_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(table_x - 0.1), Inches(table_y - 0.1),
            Inches(table_width + 0.2), Inches(num_rows * row_height + 0.2)
        )
        table_bg.fill.solid()
        table_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        table_bg.fill.fore_color.brightness = 0.3
        table_bg.line.color.rgb = self.colors['primary']
        table_bg.line.width = Pt(2)

        # 添加表格内容
        for row_idx, row in enumerate(table_data):
            for col_idx, cell_text in enumerate(row):
                # 计算单元格位置
                cell_x = table_x + col_idx * col_width
                cell_y = table_y + row_idx * row_height

                # 判断是否为表头
                is_header = (row_idx == 0)

                # 单元格背景
                cell_bg = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(cell_x), Inches(cell_y),
                    Inches(col_width), Inches(row_height)
                )
                cell_bg.fill.solid()
                if is_header:
                    cell_bg.fill.fore_color.rgb = self.colors['primary']
                    cell_bg.fill.fore_color.brightness = 0.3
                else:
                    cell_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    cell_bg.fill.fore_color.brightness = 0.4
                cell_bg.line.color.rgb = self.colors['primary']
                cell_bg.line.width = Pt(1)

                # 单元格文本
                text_box = slide.shapes.add_textbox(
                    Inches(cell_x + 0.1), Inches(cell_y + 0.05),
                    Inches(col_width - 0.2), Inches(row_height - 0.1)
                )
                text_frame = text_box.text_frame
                text_frame.text = cell_text
                text_frame.word_wrap = True
                para = text_frame.paragraphs[0]
                para.alignment = PP_ALIGN.LEFT
                para.font.size = Pt(16 if is_header else 14)
                para.font.bold = is_header
                para.font.color.rgb = self.colors['text_primary'] if is_header else self.colors['text_secondary']
                para.font.name = self.fonts.get('body', '微软雅黑')

    def add_content_slide(self, title: str, items: List[str]) -> None:
        """添加内容幻灯片（列表）"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 添加背景
        self.add_background(slide)

        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.8), Inches(8), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['text_primary']
        title_para.font.name = self.fonts.get('heading', '微软雅黑')

        # 添加列表项
        start_y = 2.2
        for item in items:
            # 创建列表项背景框
            item_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.5), Inches(start_y),
                Inches(7), Inches(0.7)
            )
            item_box.fill.solid()
            item_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
            item_box.fill.fore_color.brightness = 0.4
            item_box.line.color.rgb = self.colors['primary']
            item_box.line.width = Pt(2)

            # 添加列表项文本
            text_box = slide.shapes.add_textbox(
                Inches(1.7), Inches(start_y + 0.1),
                Inches(6.6), Inches(0.5)
            )
            text_frame = text_box.text_frame
            text_frame.text = f"• {item}"
            text_frame.word_wrap = True
            # 启用自动调整大小
            text_frame.auto_size = True
            para = text_frame.paragraphs[0]
            para.font.size = Pt(20)  # 减小字体以适应更多内容
            para.font.color.rgb = self.colors['text_secondary']
            para.font.name = self.fonts.get('body', '微软雅黑')

            start_y += 0.9

    def add_color_cards_slide(self, title: str, colors: List[tuple]) -> None:
        """添加色彩卡片幻灯片"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 添加背景
        self.add_background(slide)

        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(8), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['text_primary']
        title_para.font.name = self.fonts.get('heading', '微软雅黑')

        # 添加色彩卡片（2x2网格）
        positions = [
            (1.5, 1.8), (5.5, 1.8),
            (1.5, 4.2), (5.5, 4.2)
        ]

        for i, (color_name, color_meaning, color_rgb) in enumerate(colors):
            x, y = positions[i]

            # 卡片背景
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y),
                Inches(3), Inches(2)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.fill.fore_color.brightness = 0.2
            card.line.color.rgb = self.colors['primary']
            card.line.width = Pt(1)

            # 色彩圆圈
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x + 1), Inches(y + 0.3),
                Inches(1), Inches(1)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = RGBColor(*color_rgb)
            circle.line.color.rgb = RGBColor(0, 0, 0)
            circle.line.width = Pt(1)

            # 色彩名称
            name_box = slide.shapes.add_textbox(
                Inches(x), Inches(y + 1.4),
                Inches(3), Inches(0.3)
            )
            name_frame = name_box.text_frame
            name_frame.text = color_name
            name_para = name_frame.paragraphs[0]
            name_para.alignment = PP_ALIGN.CENTER
            name_para.font.size = Pt(18)
            name_para.font.bold = True
            name_para.font.color.rgb = self.colors['text_primary']
            name_para.font.name = self.fonts.get('body', '微软雅黑')

            # 色彩含义
            meaning_box = slide.shapes.add_textbox(
                Inches(x), Inches(y + 1.7),
                Inches(3), Inches(0.25)
            )
            meaning_frame = meaning_box.text_frame
            meaning_frame.text = color_meaning
            meaning_para = meaning_frame.paragraphs[0]
            meaning_para.alignment = PP_ALIGN.CENTER
            meaning_para.font.size = Pt(14)
            meaning_para.font.color.rgb = self.colors['text_light']
            meaning_para.font.name = self.fonts.get('body', '微软雅黑')

    def add_poems_slide(self, title: str, poems: List[Tuple[str, str]]) -> None:
        """添加诗句幻灯片"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 添加背景
        self.add_background(slide)

        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.3), Inches(8), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['text_primary']
        title_para.font.name = self.fonts.get('heading', '微软雅黑')

        # 添加诗句卡片
        start_y = 1.3
        for poem_text, poem_author in poems:
            # 诗句卡片背景
            poem_card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.5), Inches(start_y),
                Inches(7), Inches(1.5)
            )
            poem_card.fill.solid()
            poem_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            poem_card.fill.fore_color.brightness = 0.3
            poem_card.line.color.rgb = self.colors['primary']
            poem_card.line.width = Pt(2)

            # 诗句文本
            poem_box = slide.shapes.add_textbox(
                Inches(1.7), Inches(start_y + 0.3),
                Inches(6.6), Inches(0.8)
            )
            poem_frame = poem_box.text_frame
            poem_frame.text = poem_text
            poem_frame.word_wrap = True
            poem_para = poem_frame.paragraphs[0]
            poem_para.alignment = PP_ALIGN.CENTER
            poem_para.font.size = Pt(22)
            poem_para.font.color.rgb = self.colors['text_primary']
            poem_para.font.name = self.fonts.get('poem', '楷体')

            # 作者
            author_box = slide.shapes.add_textbox(
                Inches(1.7), Inches(start_y + 1.1),
                Inches(6.6), Inches(0.3)
            )
            author_frame = author_box.text_frame
            author_frame.text = poem_author
            author_para = author_frame.paragraphs[0]
            author_para.alignment = PP_ALIGN.CENTER
            author_para.font.size = Pt(14)
            author_para.italic = True
            author_para.font.color.rgb = self.colors['text_light']
            author_para.font.name = self.fonts.get('body', '微软雅黑')

            start_y += 1.7

    def add_two_column_slide(self, title: str, two_column: Dict[str, Any]) -> None:
        """添加两列布局幻灯片"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 添加背景
        self.add_background(slide)

        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.8), Inches(8), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['text_primary']
        title_para.font.name = self.fonts.get('heading', '微软雅黑')

        # 左列
        left_data = two_column.get('left', {})
        left_title = left_data.get('title', '')
        left_items = left_data.get('items', [])

        # 左列标题
        if left_title:
            left_title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2.0),
                Inches(4), Inches(0.5)
            )
            left_title_frame = left_title_box.text_frame
            left_title_frame.text = left_title
            left_title_para = left_title_frame.paragraphs[0]
            left_title_para.font.size = Pt(22)
            left_title_para.font.bold = True
            left_title_para.font.color.rgb = self.colors['primary']
            left_title_para.font.name = self.fonts.get('body', '微软雅黑')

        # 左列列表项
        left_start_y = 2.6
        for item in left_items:
            text_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(left_start_y),
                Inches(3.6), Inches(0.4)
            )
            text_frame = text_box.text_frame
            text_frame.text = f"• {item}"
            text_frame.word_wrap = True
            text_frame.auto_size = True
            para = text_frame.paragraphs[0]
            para.font.size = Pt(16)
            para.font.color.rgb = self.colors['text_secondary']
            para.font.name = self.fonts.get('body', '微软雅黑')
            left_start_y += 0.5

        # 右列
        right_data = two_column.get('right', {})
        right_title = right_data.get('title', '')
        right_items = right_data.get('items', [])

        # 右列标题
        if right_title:
            right_title_box = slide.shapes.add_textbox(
                Inches(5.5), Inches(2.0),
                Inches(4), Inches(0.5)
            )
            right_title_frame = right_title_box.text_frame
            right_title_frame.text = right_title
            right_title_para = right_title_frame.paragraphs[0]
            right_title_para.font.size = Pt(22)
            right_title_para.font.bold = True
            right_title_para.font.color.rgb = self.colors['primary']
            right_title_para.font.name = self.fonts.get('body', '微软雅黑')

        # 右列列表项
        right_start_y = 2.6
        for item in right_items:
            text_box = slide.shapes.add_textbox(
                Inches(5.7), Inches(right_start_y),
                Inches(3.6), Inches(0.4)
            )
            text_frame = text_box.text_frame
            text_frame.text = f"• {item}"
            text_frame.word_wrap = True
            text_frame.auto_size = True
            para = text_frame.paragraphs[0]
            para.font.size = Pt(16)
            para.font.color.rgb = self.colors['text_secondary']
            para.font.name = self.fonts.get('body', '微软雅黑')
            right_start_y += 0.5

    def add_cards_slide(self, title: str, cards: List[Dict[str, str]]) -> None:
        """添加卡片幻灯片"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 添加背景
        self.add_background(slide)

        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(8), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['text_primary']
        title_para.font.name = self.fonts.get('heading', '微软雅黑')

        # 检查是否只有一个大卡片（如"总结"或"故事背景"）
        if len(cards) == 1:
            # 单个大卡片布局
            card = cards[0]
            x, y = 1.5, 1.8
            width, height = 7.0, 4.5

            # 卡片背景
            card_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y),
                Inches(width), Inches(height)
            )
            card_shape.fill.solid()
            card_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card_shape.fill.fore_color.brightness = 0.2
            card_shape.line.color.rgb = self.colors['primary']
            card_shape.line.width = Pt(1)

            # 卡片标题（如果有）
            card_title = card.get('title', '')
            if card_title:
                title_box = slide.shapes.add_textbox(
                    Inches(x), Inches(y + 0.3),
                    Inches(width), Inches(0.5)
                )
                title_frame = title_box.text_frame
                title_frame.text = card_title
                title_para = title_frame.paragraphs[0]
                title_para.alignment = PP_ALIGN.CENTER
                title_para.font.size = Pt(24)
                title_para.font.bold = True
                title_para.font.color.rgb = self.colors['text_primary']
                title_para.font.name = self.fonts.get('body', '微软雅黑')
                content_start_y = y + 1.0
            else:
                # 没有标题，内容从顶部开始
                content_start_y = y + 0.5

            # 卡片内容（支持多段落）
            content_text = card.get('content', '')
            paragraphs = content_text.split('\n')

            content_box = slide.shapes.add_textbox(
                Inches(x + 0.3), Inches(content_start_y),
                Inches(width - 0.6), Inches(height - (content_start_y - y) - 0.3)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            # 添加每个段落
            for i, para_text in enumerate(paragraphs):
                if i == 0:
                    para = content_frame.paragraphs[0]
                else:
                    para = content_frame.add_paragraph()
                para.text = para_text
                para.font.size = Pt(18)
                para.font.color.rgb = self.colors['text_secondary']
                para.font.name = self.fonts.get('body', '微软雅黑')
                para.space_after = Pt(12)
        else:
            # 多个小卡片布局（2x2网格）
            positions = [
                (1.5, 1.8), (5.5, 1.8),
                (1.5, 4.2), (5.5, 4.2)
            ]

            for i, card in enumerate(cards):
                if i >= len(positions):
                    break

                x, y = positions[i]

                # 卡片背景
                card_shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(x), Inches(y),
                    Inches(3), Inches(2)
                )
                card_shape.fill.solid()
                card_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
                card_shape.fill.fore_color.brightness = 0.2
                card_shape.line.color.rgb = self.colors['primary']
                card_shape.line.width = Pt(1)

                # 卡片标题
                title_box = slide.shapes.add_textbox(
                    Inches(x), Inches(y + 0.3),
                    Inches(3), Inches(0.5)
                )
                title_frame = title_box.text_frame
                title_frame.text = card.get('title', '')
                title_para = title_frame.paragraphs[0]
                title_para.alignment = PP_ALIGN.CENTER
                title_para.font.size = Pt(18)
                title_para.font.bold = True
                title_para.font.color.rgb = self.colors['text_primary']
                title_para.font.name = self.fonts.get('body', '微软雅黑')

                # 卡片内容
                content_box = slide.shapes.add_textbox(
                    Inches(x + 0.2), Inches(y + 0.9),
                    Inches(2.6), Inches(0.9)
                )
                content_frame = content_box.text_frame
                content_frame.text = card.get('content', '')
                content_frame.word_wrap = True
                content_para = content_frame.paragraphs[0]
                content_para.font.size = Pt(14)
                content_para.font.color.rgb = self.colors['text_secondary']
                content_para.font.name = self.fonts.get('body', '微软雅黑')

    def add_image_slide(self, title: str, image_data: Dict[str, str]) -> None:
        """添加图片幻灯片"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 添加背景
        self.add_background(slide)

        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(8), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['text_primary']
        title_para.font.name = self.fonts.get('heading', '微软雅黑')

        # 添加图片
        image_path = image_data.get('src', '')
        if image_path and os.path.exists(image_path):
            slide.shapes.add_picture(
                image_path,
                Inches(2), Inches(1.8),
                width=Inches(6)
            )

        # 添加说明文字
        caption = image_data.get('caption', '')
        if caption:
            caption_box = slide.shapes.add_textbox(
                Inches(1), Inches(6),
                Inches(8), Inches(0.5)
            )
            caption_frame = caption_box.text_frame
            caption_frame.text = caption
            caption_para = caption_frame.paragraphs[0]
            caption_para.alignment = PP_ALIGN.CENTER
            caption_para.font.size = Pt(16)
            caption_para.font.color.rgb = self.colors['text_secondary']
            caption_para.font.name = self.fonts.get('body', '微软雅黑')

    def add_decoration(self, slide) -> None:
        """添加装饰性元素"""
        positions = [
            (0.3, 0.3, 0.3), (8.7, 0.5, 0.25),
            (0.5, 6.5, 0.25), (8.5, 6.3, 0.3)
        ]

        for x, y, size in positions:
            flower = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x), Inches(y),
                Inches(size), Inches(size)
            )
            flower.fill.solid()
            flower.fill.fore_color.rgb = self.colors['secondary']
            flower.fill.fore_color.brightness = 0.3
            flower.line.fill.background()

    def verify_editability(self) -> Dict[str, int]:
        """验证PPTX可编辑性"""
        stats = {
            'text_boxes': 0,
            'shapes': 0,
            'images': 0,
            'screenshots': 0
        }

        for slide in self.prs.slides:
            for shape in slide.shapes:
                stats['shapes'] += 1
                if shape.has_text_frame:
                    stats['text_boxes'] += 1
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    stats['images'] += 1

        return stats

    def save(self, output_file: str) -> None:
        """保存PPTX文件"""
        self.prs.save(output_file)
